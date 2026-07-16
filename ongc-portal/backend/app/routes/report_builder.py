from fastapi import APIRouter, Depends, HTTPException, Form
from fastapi.responses import Response
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from sqlalchemy.orm import selectinload
from app.database import get_db
from app.models.base import ReportTemplate, ReportPeriod, ReportSubmission, User, Notification
from app.auth.deps import get_current_user
from datetime import datetime, date
import json

router = APIRouter()

# ─── Helper ───
def admin_or_ops(user):
    role = user.role.name if user.role else ""
    if role not in ("admin", "ops_manager"):
        raise HTTPException(403, "Admin or Ops Manager access required")

def any_auth(user):
    if not user:
        raise HTTPException(401, "Authentication required")

# ─── Templates ───

@router.get("/templates")
async def list_templates(db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    any_auth(user)
    result = await db.execute(
        select(ReportTemplate).options(selectinload(ReportTemplate.creator)).order_by(ReportTemplate.created_at.desc())
    )
    templates = result.scalars().all()
    return [{
        "id": t.id,
        "name": t.name,
        "description": t.description,
        "period_type": t.period_type,
        "sections": json.loads(t.sections) if isinstance(t.sections, str) else t.sections,
        "assigned_roles": json.loads(t.assigned_roles) if isinstance(t.assigned_roles, str) else (t.assigned_roles or []),
        "section": t.section,
        "area": t.area,
        "created_by_name": t.creator.name if t.creator else None,
        "created_at": str(t.created_at) if t.created_at else None,
    } for t in templates]

@router.post("/templates/create")
async def create_template(
    name: str = Form(...),
    description: str = Form(None),
    period_type: str = Form("monthly"),
    sections: str = Form(...),
    assigned_roles: str = Form("[]"),
    section: str = Form(None),
    area: str = Form(None),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    if not name or not sections:
        raise HTTPException(400, "Name and sections required")
    try:
        json.loads(sections)
    except:
        raise HTTPException(400, "Sections must be valid JSON")
    try:
        roles = json.loads(assigned_roles)
        if not isinstance(roles, list):
            raise ValueError
    except:
        raise HTTPException(400, "assigned_roles must be a JSON array")
    t = ReportTemplate(
        name=name,
        description=description,
        period_type=period_type,
        sections=sections,
        assigned_roles=json.dumps(roles),
        section=section,
        area=area,
        created_by=user.id,
    )
    db.add(t)
    await db.commit()
    await db.refresh(t)
    return {"id": t.id, "name": t.name, "msg": "Template created"}

@router.put("/templates/{template_id}")
async def update_template(
    template_id: int,
    name: str = Form(None),
    description: str = Form(None),
    period_type: str = Form(None),
    sections: str = Form(None),
    assigned_roles: str = Form(None),
    section: str = Form(None),
    area: str = Form(None),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    result = await db.execute(select(ReportTemplate).where(ReportTemplate.id == template_id))
    t = result.scalar_one_or_none()
    if not t:
        raise HTTPException(404, "Template not found")
    if name: t.name = name
    if description is not None: t.description = description
    if period_type: t.period_type = period_type
    if sections:
        try:
            json.loads(sections)
        except:
            raise HTTPException(400, "Sections must be valid JSON")
        t.sections = sections
    if assigned_roles:
        try:
            roles = json.loads(assigned_roles)
            if not isinstance(roles, list):
                raise ValueError
            t.assigned_roles = json.dumps(roles)
        except:
            raise HTTPException(400, "assigned_roles must be a JSON array")
    if section is not None: t.section = section
    if area is not None: t.area = area
    await db.commit()
    return {"msg": "Template updated"}

@router.delete("/templates/{template_id}")
async def delete_template(template_id: int, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    admin_or_ops(user)
    result = await db.execute(select(ReportTemplate).where(ReportTemplate.id == template_id))
    t = result.scalar_one_or_none()
    if not t:
        raise HTTPException(404, "Template not found")
    await db.delete(t)
    await db.commit()
    return {"msg": "Template deleted"}

# ─── Periods ───

@router.get("/periods")
async def list_periods(template_id: int = None, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    any_auth(user)
    q = select(ReportPeriod).order_by(ReportPeriod.created_at.desc())
    if template_id:
        q = q.where(ReportPeriod.template_id == template_id)
    result = await db.execute(q)
    periods = result.scalars().all()
    return [{
        "id": p.id,
        "template_id": p.template_id,
        "label": p.label,
        "start_date": str(p.start_date) if p.start_date else None,
        "end_date": str(p.end_date) if p.end_date else None,
        "is_open": p.is_open,
        "section_assignments": json.loads(p.section_assignments) if p.section_assignments else {},
        "created_at": str(p.created_at) if p.created_at else None,
    } for p in periods]

@router.post("/periods/create")
async def create_period(
    template_id: int = Form(...),
    label: str = Form(...),
    start_date: str = Form(None),
    end_date: str = Form(None),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    p = ReportPeriod(
        template_id=template_id,
        label=label,
        start_date=datetime.strptime(start_date, "%Y-%m-%d").date() if start_date else None,
        end_date=datetime.strptime(end_date, "%Y-%m-%d").date() if end_date else None,
    )
    db.add(p)
    await db.commit()
    await db.refresh(p)
    return {"id": p.id, "label": p.label, "msg": "Period created"}

@router.post("/periods/{period_id}/close")
async def close_period(period_id: int, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    admin_or_ops(user)
    result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(404, "Period not found")
    p.is_open = False
    await db.commit()
    return {"msg": "Period closed"}

@router.put("/periods/{period_id}/assignments")
async def update_assignments(
    period_id: int,
    assignments: str = Form(...),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    admin_or_ops(user)
    try:
        parsed = json.loads(assignments)
    except:
        raise HTTPException(400, "Invalid JSON")
    result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(404, "Period not found")
    p.section_assignments = json.dumps(parsed)
    await db.commit()
    return {"msg": "Assignments updated"}

@router.post("/seed-demo")
async def seed_demo_data(db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    admin_or_ops(user)
    import random as _random

    existing_periods = await db.execute(select(ReportPeriod))
    existing_count = len(existing_periods.scalars().all())

    result = await db.execute(select(ReportTemplate))
    templates = result.scalars().all()
    if not templates:
        raise HTTPException(400, "No templates found — create templates first")

    result = await db.execute(select(User).where(User.is_active == True))
    all_users = result.scalars().all()
    demo_cpfs = [
        "100001","100002","100003","100006","100007","100008",
        "100010","100012","100013","100014","100015","100018","100019",
    ]
    demo_users = [u for u in all_users if u.cpf in demo_cpfs]
    if not demo_users:
        demo_users = all_users[:13]

    RIG_NAMES = [
        "ONGC Rig-1","ONGC Rig-3","ONGC Rig-7","ONGC Rig-12","ONGC Rig-15",
        "Samsung Rig S-201","Transocean DDK-1","Aban Pearl","Bela Driller",
        "Jindal Star","GreatDrill R-4","Essar E-12","Adani Jackup J-1",
    ]
    WELL_NAMES = [
        "CB-1/1A","CB-1/2","CB-2/1","CB-3/1","CB-3/2","CB-4/1",
        "MBS-1","MBS-2","MBS-3","D-1","D-2","D-3","D-4",
        "KG-9","KG-10","KG-11","KG-12","PR-1","PR-2","PR-3",
    ]
    LOCATIONS = [
        "Mumbai High North","Mumbai High South","Neelam Field","Heera Platform",
        "Panna-Mukta","Bombay High","Rajahmundry","Kakinada",
        "Ahmedabad Asset","Jodhpur","Dehradun","Kolkata",
        "Cauvery Basin","Krishna-Godavari Basin","Assam Asset","Cambay Basin",
    ]
    EQUIPMENT = [
        "Topside Process Module","Gas Compressor Unit","Water Injection Pump",
        "Crude Oil Transfer Pump","Emergency Generator","Fire Water System",
        "BOP Stack","Mud Circulating System","Derrick Substructure",
        "Draw Works","Rotary Table","Casing Running Tool","Wellhead Assembly",
        "Christmas Tree","Separator Unit","Desalter","Heat Exchanger",
    ]
    ACTIVITIES = [
        "Drilling operations in progress","Workover activity on well",
        "Well completion and testing","Maintenance shutdown for inspection",
        "Platform turnaround maintenance","Pipeline pigging operation",
        "Seismic survey acquisition","Well stimulation activity",
        " artificial lift optimization","Produced water treatment",
        "Gas lift valve replacement","ESP installation and commissioning",
        "Wellhead pressure testing","Flowline leak repair",
    ]
    CHALLENGES = [
        "Monsoon season affecting offshore logistics and crew change schedules",
        "Equipment delivery delays due to port congestion at Nhava Sheva",
        "Higher than expected sand production requiring frequent well intervention",
        "Crew fatigue due to extended rotation schedules during peak operations",
        "Budget constraints limiting procurement of critical spare parts",
        "Weather window limitations for jack-up rig positioning",
        "Downhole complications requiring sidetrack drilling",
        "Regulatory inspection timelines creating operational pressure",
    ]
    RECOMMENDATIONS = [
        "Increase standby equipment inventory to reduce unplanned downtime",
        "Implement predictive maintenance for rotating equipment using vibration analysis",
        "Optimize crew rotation schedules to improve productivity and safety",
        "Accelerate procurement process for long-lead items through framework agreements",
        "Deploy additional PSVs for offshore logistics during monsoon season",
        "Conduct weekly safety stand-down meetings to reinforce HSE protocols",
        "Invest in digital monitoring systems for real-time well performance tracking",
        "Review and update emergency response plans for all offshore installations",
    ]
    TECH_ADOPTIONS = [
        "AI-based predictive analytics for equipment failure detection",
        "Digital twin simulation for reservoir management optimization",
        "Drone-based inspection of flare stacks and elevated structures",
        "Remote operations center enabling real-time decision making",
        "Advanced LWD tools for geosteering in thin reservoir sections",
        "Smart well completion technology for improved oil recovery",
    ]
    STRATEGIC_INITIATIVES = [
        "Digital transformation of field operations through IoT deployment",
        "Green hydrogen pilot project integration with existing facilities",
        "Enhanced oil recovery program using polymer flooding technique",
        "Subsea tieback development for marginal field exploitation",
        "Carbon capture and storage feasibility study for mature fields",
        "Workforce upskilling program for digital oilfield technologies",
    ]
    KPI_NAMES = [
        "Oil Production Efficiency","Gas Recovery Rate","Water Cut Reduction",
        "Equipment Reliability Index","Drilling Penetration Rate",
        "Workover Success Rate","HSE Compliance Score","Cost per Barrel",
    ]

    def _rand_month_name(m):
        return ["","January","February","March","April","May","June",
                "July","August","September","October","November","December"][m]

    def _fy_label(year, month):
        if month >= 4:
            return f"FY {year}-{str(year+1)[-2:]}"
        return f"FY {year-1}-{str(year)[-2:]}"

    def gen_field_value(field, sec_key, user, month, year, rng):
        ftype = field.get("type", "text")
        key = field.get("key", "")
        label = field.get("label", key)

        if ftype == "number":
            if "progress" in key or "pct" in key or "achievement" in key:
                return str(rng.randint(45, 98))
            if "depth" in key and "target" in key:
                return str(rng.randint(2000, 4500))
            if "depth" in key:
                return str(rng.randint(800, 3800))
            if "downtime" in key:
                return str(round(rng.uniform(0.5, 48.0), 1))
            if "manhours" in key or "lti" in key.lower():
                return str(rng.randint(8000, 25000))
            if "staff" in key or "headcount" in key:
                base = rng.randint(120, 350)
                if "field" in key:
                    return str(int(base * rng.uniform(0.5, 0.7)))
                if "office" in key:
                    return str(int(base * rng.uniform(0.3, 0.5)))
                return str(base)
            if "vacanc" in key:
                return str(rng.randint(2, 18))
            if "budget" in key or "allocated" in key or "utilized" in key or "expenditure" in key or "remaining" in key:
                base = rng.randint(500, 3000)
                if "utilized" in key or "expenditure" in key:
                    return str(base + rng.randint(-200, 300))
                if "remaining" in key:
                    return str(max(50, base - rng.randint(100, 600)))
                return str(base)
            if "volume" in key or "skm" in label.lower():
                return str(round(rng.uniform(12.5, 85.0), 1))
            if "project" in key and "complet" in key:
                return str(rng.randint(3, 15))
            if "project" in key and "ongoing" in key:
                return str(rng.randint(5, 20))
            if "project" in key and "planned" in key:
                return str(rng.randint(8, 25))
            if "kpi" in key and "overall" in key:
                return str(rng.randint(72, 96))
            if "savings" in key.lower():
                return f"₹{rng.randint(2, 45)} Cr {'savings' if rng.random() > 0.3 else 'overrun'}"
            return str(rng.randint(1, 500))

        if ftype == "select":
            if "status" in key or "overall" in key:
                return rng.choice(["On Track","At Risk","Completed","Delayed","In Progress","Operational","Under Maintenance"])
            return rng.choice(["On Track","Completed","Operational","In Progress","At Risk","Delayed"])

        if ftype == "date":
            return f"{year}-{month:02d}-{rng.randint(1,28):02d}"

        if ftype == "textarea" or ftype == "text":
            sec = sec_key.lower()

            if sec in ("ops_update", "ops"):
                if "activity" in key:
                    return rng.choice(ACTIVITIES)
                if "location" in key or "rig" in key.lower():
                    return rng.choice(RIG_NAMES) + " — " + rng.choice(LOCATIONS)
                if "planned" in key or "plan" in key:
                    plans = [
                        "Continue drilling to next casing point at 2800m TD. Complete DST on Zone-3. Mobilize wireline unit for log acquisition.",
                        "Commence well completion activities. Install ESP and surface facilities. Begin flowback testing within 72 hours.",
                        "Proceed with platform turnaround maintenance. Replace corroded flowline sections. Conduct hydrostatic testing on new lines.",
                        "Begin workover operations on well. Pull completion string. Run diagnostic logs and evaluate reservoir conditions.",
                        "Complete pipeline pigging operations. Inspect cathodic protection system. Schedule NDT inspection for critical welds.",
                    ]
                    return rng.choice(plans)
                if "progress" in key:
                    return str(rng.randint(40, 95))
                if "key" in key.lower() or "highlight" in key.lower():
                    highlights = [
                        "Achieved record drilling footage of 320m in 24 hours. Zero LTI maintained for 450 days. New well spudded ahead of schedule.",
                        "Successfully completed well stimulation increasing production by 18%. Installed new separator unit reducing processing downtime.",
                        "Monsoon preparedness checklist 100% complete. Safety audit passed with zero non-conformities. Crew training program completed.",
                        "ESP replacement completed in 48 hours minimizing production loss. Pipeline integrity assessment finished with no critical findings.",
                        "Quarterly target achieved 92% with improved efficiency metrics. Equipment reliability index improved to 97.3% from 94.1%.",
                    ]
                    return rng.choice(highlights)
                if "summary" in key.lower():
                    summaries = [
                        "Operations running smoothly with minor weather-related delays. All critical equipment operational. Safety performance excellent with zero incidents this period.",
                        "Production targets met despite planned maintenance window. Drilling operations progressing on schedule. New well shows promising initial production rates.",
                        "Mixed performance with one equipment failure causing 12-hour production dip. Corrective actions implemented. Overall field performance above expectation.",
                    ]
                    return rng.choice(summaries)
                return rng.choice([
                    "Routine operations ongoing. All systems nominal. No safety concerns reported.",
                    "Active drilling with positive reservoir indications. Equipment performance within parameters.",
                    "Maintenance activities completed ahead of schedule. Ready for next phase operations.",
                ])

            if sec in ("drilling",):
                if "well" in key:
                    return rng.choice(WELL_NAMES)
                if "remark" in key:
                    remarks = [
                        "Drilling ahead at 15m/hr through tertiary sandstone. Good ROP maintained. Mud weight optimized for formation pressure.",
                        "Encountered lost circulation zone at 2450m. Applied LCM treatment successfully. Resumed drilling after 6-hour cure period.",
                        "Completed DST on Zone-2. Flow rate: 1200 bopd. Oil sample collected for PVT analysis. Well declared commercially viable.",
                        "Casing cement job completed successfully. Top of cement verified at surface. Pressure test passed at 2500 psi for 30 minutes.",
                        "Sidetrack drilling initiated due to stuck pipe incident. Current depth 2100m. Wellbore geometry satisfactory for continued operations.",
                    ]
                    return rng.choice(remarks)
                return rng.randint(800, 4200)

            if sec in ("equipment",):
                if "equip" in key:
                    return rng.choice(EQUIPMENT)
                if "status" in key:
                    return rng.choice(["Operational","Under Maintenance","Completed","In Progress","At Risk"])
                if "down" in key:
                    return str(round(rng.uniform(0.5, 36.0), 1))
                return rng.choice(EQUIPMENT)

            if sec in ("hse",):
                if "incident" in key:
                    incidents = [
                        "No reportable incidents this period. Minor slip incident during crew change — first aid administered, no lost time.",
                        "One recordable incident: finger injury during pipe handling. Worker treated at field clinic and returned to duty same day.",
                        "Near-miss reported: dropped object from crane during lift operation. Area cordoned off, investigation completed within 24 hours.",
                        "Zero incidents recorded. 2 near-miss reports submitted and investigated. All corrective actions closed within stipulated time.",
                        "Minor spill during fuel transfer — contained on deck, no environmental impact. Procedure reviewed and crew re-briefed.",
                    ]
                    return rng.choice(incidents)
                if "safety" in key.lower() or "observation" in key.lower():
                    observations = [
                        "12 safety observations logged: 8 positive, 4 corrective. Focus areas: PPE compliance, working at height, confined space entry.",
                        "25 safety observations recorded during monthly walkthrough. All high-risk activities covered with concurrent safety monitoring.",
                        "9 observations submitted — 6 positive reinforcement, 3 requiring improvement. Toolbox talks conducted for all high-risk activities.",
                        "15 observations covering LOTO procedures, scaffold inspections, and crane operations. All critical observations addressed immediately.",
                    ]
                    return rng.choice(observations)
                return str(rng.randint(8000, 22000))

            if sec in ("issues",):
                if "issue" in key:
                    issue_texts = [
                        "Corroded flowline section identified during pigging run — requires replacement within 30 days to prevent leak.",
                        "Gas compressor vibration levels exceeding threshold — bearing replacement scheduled for next maintenance window.",
                        "Crude oil pump seal failure causing intermittent production stoppage — spare seal kit on order, expected delivery in 5 days.",
                        "Regulatory audit findings pending closure — 3 observations require documentation updates by end of month.",
                        "Crew shortage for night shift operations — request for additional manpower from Kolkata office pending approval.",
                    ]
                    return rng.choice(issue_texts)
                if "resolution" in key or "action" in key:
                    resolutions = [
                        "Emergency procurement initiated for replacement flowline. Temporary clamping applied. Monitoring daily for any deterioration.",
                        "Bearing replacement completed during planned 48-hour shutdown. Vibration levels now within acceptable limits. Continuous monitoring enabled.",
                        "Seal kit received and installed. Pump restored to full capacity. Root cause analysis indicates contamination — inline filter added.",
                        "Documentation updated and submitted to regulatory authority. Internal review process enhanced to prevent recurrence.",
                        "Temporary arrangement with contractor for night shift coverage. Permanent hiring proposal submitted to HR for 4 additional operators.",
                    ]
                    return rng.choice(resolutions)
                if "escalation" in key:
                    return rng.choice(["No — resolved at field level","Yes — escalated to Asset Manager","No — action plan approved","Yes — pending management review"])
                return rng.choice([
                    "No critical issues reported. All ongoing items being tracked through standard workflow.",
                    "Two items requiring management attention — budget approval and resource allocation.",
                ])

            if sec in ("exec_summary",):
                if "highlight" in key:
                    return rng.choice([
                        "Achieved 95% of monthly production target despite one planned shutdown. New well D-4 contributed 800 bopd additional production.",
                        "Safety milestone: 500 days without LTI. Drilling program on track with 2 wells completed ahead of schedule.",
                        "Budget utilization at 87% with quarterly savings of ₹12 Cr through optimized procurement. Equipment reliability improved to 97%.",
                    ])
                if "summary" in key:
                    return rng.choice([
                        "Strong operational performance with zero safety incidents. Production slightly below target due to weather-related shutdowns in second fortnight. All regulatory compliance requirements met.",
                        "Excellent month with production exceeding target by 3%. New well commissioning completed successfully. Manpower utilization improved with targeted training programs.",
                        "Mixed performance with equipment availability at 94%. Two critical spares received from international suppliers. Budget on track with minor reallocation needed.",
                    ])
                return rng.choice(["On Track","At Risk","Completed"])

            if sec in ("project_progress",):
                if "detail" in key:
                    return rng.choice([
                        "Well D-4: Completed — production commenced at 800 bopd. Well D-5: Drilling at 2800m TD. Pipeline leak repair: 90% complete.",
                        "ESP installation on MBS-2 completed. Workover on Heera-A7 in progress. Flowline replacement project 75% complete.",
                        "Seismic data acquisition completed for Block-XY. Processing underway. Two workover jobs completed successfully.",
                    ])
                return rng.randint(1, 25)

            if sec in ("targets",):
                if "deviation" in key:
                    return rng.choice([
                        "Minor deviation due to planned maintenance shutdown lasting 3 days. Corrective drilling schedule adjusted for next period.",
                        "Over-achievement attributed to successful well stimulation on MBS-3. Additional production of 200 bopd sustained.",
                        "Weather-related production loss of 2 days. Compensatory measures through extended well testing planned.",
                    ])
                return rng.randint(10, 90)

            if sec in ("manpower",):
                return str(rng.randint(2, 18))

            if sec in ("budget",):
                if "note" in key:
                    return rng.choice([
                        "Budget on track. Major spend categories: drilling (42%), maintenance (28%), logistics (18%), admin (12%).",
                        "Minor overspend on emergency repairs offset by underspend on planned maintenance. Net variance within acceptable limits.",
                        "Capital expenditure 78% utilized. Remaining allocation for Q4 equipment replacement program.",
                    ])
                return f"₹{rng.randint(100, 5000)} Cr"

            if sec in ("challenges",):
                if "challeng" in key:
                    return rng.choice(CHALLENGES)
                if "recommend" in key:
                    return rng.choice(RECOMMENDATIONS)
                return rng.choice(CHALLENGES)

            if sec in ("overview", "quarter"):
                if "quarter" in key:
                    q = ((month - 1) // 3) + 1
                    return f"Q{q} {_fy_label(year, month)}"
                if "achievement" in key or "metric" in key:
                    return rng.choice([
                        "Record production achieved in Month 12. Two new wells commissioned ahead of schedule. Safety performance exemplary with zero LTIs.",
                        "Drilling program completed 85% of annual target. Budget utilization efficient at 91%. Manpower productivity improved by 12%.",
                        "All major milestones achieved. Revenue target exceeded by 5%. Operational expenditure maintained within approved budget.",
                    ])
                return f"Q{((month-1)//3)+1} {_fy_label(year, month)}"

            if sec in ("kpi",):
                if "kpi" in key and "overall" not in key:
                    kpi = rng.choice(KPI_NAMES)
                    val = rng.randint(78, 99)
                    return f"{kpi}: {val}%"
                return str(rng.randint(75, 98))

            if sec in ("technical",):
                if "tech" in key or "adopt" in key:
                    return rng.choice(TECH_ADOPTIONS)
                if "innovat" in key:
                    return rng.choice([
                        "Developed automated well test data interpretation tool reducing analysis time by 60%.",
                        "Implemented real-time drilling optimization using machine learning algorithms — ROP improved by 15%.",
                        "Created digital dashboard for HSE metrics tracking across all offshore installations.",
                    ])
                if "publicat" in key:
                    return rng.choice([
                        "2 papers submitted to SPE for Annual Technical Conference. 1 paper accepted in Journal of Petroleum Technology.",
                        "Case study on ESP optimization published in ONGC Technical Bulletin. Conference presentation at India Oil & Gas Summit.",
                        "Research paper on enhanced oil recovery techniques submitted to SPE Asia Pacific. Workshop conducted for knowledge sharing.",
                    ])
                return rng.choice(TECH_ADOPTIONS)

            if sec in ("financial",):
                if "saving" in key or "overrun" in key:
                    return rng.choice([
                        f"Net savings of ₹{rng.randint(5, 80)} Cr due to optimized procurement and reduced downtime",
                        f"Overrun of ₹{rng.randint(2, 35)} Cr due to emergency repairs and accelerated procurement",
                        "Within budget with minor savings from operational efficiency improvements",
                    ])
                return f"₹{rng.randint(50, 500)}.{rng.randint(0,9)} Cr"

            if sec in ("team",):
                if "train" in key or "certif" in key:
                    return rng.choice([
                        "12 operators certified on new BOP control system. 8 engineers completed HUET training. Safety refresher conducted for all offshore staff.",
                        "Competency assessment completed for 45 personnel. 6 operators promoted to senior roles. Annual safety training program 80% complete.",
                        "15 crew members completed advanced first-aid training. 3 supervisors enrolled in leadership development program.",
                    ])
                if "attrit" in key or "joiner" in key:
                    return rng.choice([
                        "No attrition this quarter. 8 new trainees inducted into drilling operations. 2 lateral hires for specialized maintenance roles.",
                        "1 resignation in admin section — replacement hired. 5 fresh graduates joined through campus recruitment. Net headcount +4.",
                        "Zero attrition in technical cadre. 3 contractual staff terminated due to performance. New contracts awarded to 2 vendors.",
                    ])
                return str(rng.randint(150, 400))

            if sec in ("road_ahead", "road"):
                if "plan" in key:
                    return rng.choice([
                        "Focus on drilling 3 new development wells. Commission water injection project. Target production increase of 8% over current quarter.",
                        "Complete ESP replacement program for 5 wells. Initiate EOR pilot project. Begin infrastructure upgrade for gas processing.",
                        "Accelerate drilling campaign with 2 rigs operational. Complete pipeline integrity assessment. Launch digital oilfield initiative.",
                    ])
                if "strateg" in key:
                    return rng.choice(STRATEGIC_INITIATIVES)
                if "risk" in key:
                    return rng.choice([
                        "Monsoon season (Jun-Sep) may impact offshore operations. Mitigation: advance inventory stocking and weather monitoring.",
                        "Regulatory changes in environmental compliance may require facility upgrades. Budget provision of ₹50 Cr earmarked.",
                        "Supply chain disruption risk for critical equipment — alternate vendor qualification in progress.",
                    ])
                return rng.choice(STRATEGIC_INITIATIVES)

            # fallback
            return f"Report data for {label} — period {month}/{year}, submitted by {user.name}"

        return f"{label}: value"

    # ─── Generate periods across FY 2023-24, FY 2024-25, FY 2025-26 ───
    # April 2023 to March 2026 = 36 months
    years_months = []
    for yr in [2023, 2024, 2025]:
        for m in range(4, 13):
            years_months.append((yr, m))
        for m in range(1, 4):
            years_months.append((yr + 1, m))
    # Deduplicate and sort
    years_months = sorted(set(years_months))
    # Limit to 36 months (3 fiscal years)
    years_months = years_months[:36]

    periods_created = 0
    subs_created = 0
    templates_used = []
    users_used = [u.name for u in demo_users]

    for t in templates:
        sections_data = json.loads(t.sections) if isinstance(t.sections, str) else t.sections
        templates_used.append({
            "id": t.id,
            "name": t.name,
            "sections": len(sections_data or []),
            "fields_per_section": [len(s.get("fields", [])) for s in (sections_data or [])],
        })

        for yr, m in years_months:
            days_in_month = (date(yr, m + 1, 1) if m < 12 else date(yr + 1, 1, 1)) - date(yr, m, 1)
            fy = _fy_label(yr, m)
            label = f"{t.name} — {_rand_month_name(m)} {yr} ({fy})"

            period = ReportPeriod(
                template_id=t.id,
                label=label,
                start_date=date(yr, m, 1),
                end_date=date(yr, m, min(days_in_month.days, 28)),
                is_open=(yr >= 2025),
            )
            db.add(period)
            await db.flush()
            periods_created += 1

            for du in demo_users:
                fv = {}
                if sections_data:
                    sec = _random.choice(sections_data)
                    for f in (sec.get("fields") or []):
                        fv[f["key"]] = gen_field_value(f, sec.get("key", ""), du, m, yr, _random)
                fv_json = json.dumps(fv)
                status = "submitted" if _random.random() > 0.2 else "draft"
                sub = ReportSubmission(
                    period_id=period.id,
                    section_key="__full__",
                    user_id=du.id,
                    field_values=fv_json,
                    status=status,
                    submitted_at=datetime(yr, m, min(days_in_month.days, 15), _random.randint(8, 17), _random.randint(0, 59)) if status == "submitted" else None,
                )
                db.add(sub)
                subs_created += 1

                # Also create section-level submissions for ~2 random sections
                if sections_data and len(sections_data) > 1:
                    extra_sections = _random.sample(sections_data, min(2, len(sections_data)))
                    for sec in extra_sections:
                        sec_fv = {}
                        for f in (sec.get("fields") or []):
                            sec_fv[f["key"]] = gen_field_value(f, sec.get("key", ""), du, m, yr, _random)
                        sec_sub = ReportSubmission(
                            period_id=period.id,
                            section_key=sec.get("key", "section"),
                            user_id=du.id,
                            field_values=json.dumps(sec_fv),
                            status="submitted" if _random.random() > 0.25 else "draft",
                            submitted_at=datetime(yr, m, min(days_in_month.days, 20), _random.randint(8, 17), _random.randint(0, 59)) if _random.random() > 0.25 else None,
                        )
                        db.add(sec_sub)
                        subs_created += 1

    await db.commit()

    return {
        "msg": f"Seed complete: {periods_created} periods, {subs_created} submissions across {len(years_months)} months",
        "summary": {
            "years_range": f"{years_months[0][0]}-{years_months[-1][0]}" if years_months else "N/A",
            "months_generated": len(years_months),
            "templates_used": len(templates_used),
            "template_details": templates_used,
            "demo_users_count": len(demo_users),
            "demo_users": users_used,
            "periods_created": periods_created,
            "submissions_created": subs_created,
            "existing_periods_before": existing_count,
        },
    }

@router.get("/periods/{period_id}/assignments")
async def get_assignments(period_id: int, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    any_auth(user)
    result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(404, "Period not found")
    return json.loads(p.section_assignments) if p.section_assignments else {}

# ─── Submissions ───

@router.get("/submissions")
async def list_submissions(
    period_id: int = None,
    section_key: str = None,
    assigned_to: int = None,
    user_id: int = None,
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    any_auth(user)
    q = select(ReportSubmission).options(
        selectinload(ReportSubmission.period),
        selectinload(ReportSubmission.assignee),
        selectinload(ReportSubmission.submitter),
    ).order_by(ReportSubmission.created_at.desc())
    if period_id: q = q.where(ReportSubmission.period_id == period_id)
    if section_key: q = q.where(ReportSubmission.section_key == section_key)
    if assigned_to: q = q.where(ReportSubmission.assigned_to == assigned_to)
    if user_id: q = q.where(ReportSubmission.user_id == user_id)
    result = await db.execute(q)
    subs = result.scalars().all()
    return [{
        "id": s.id,
        "period_id": s.period_id,
        "section_key": s.section_key,
        "assigned_to": s.assigned_to,
        "assigned_to_name": s.assignee.name if s.assignee else None,
        "user_id": s.user_id,
        "user_name": s.submitter.name if s.submitter else None,
        "user_cpf": s.submitter.cpf if s.submitter else None,
        "field_values": json.loads(s.field_values) if isinstance(s.field_values, str) else s.field_values,
        "status": s.status,
        "submitted_at": str(s.submitted_at) if s.submitted_at else None,
        "created_at": str(s.created_at) if s.created_at else None,
    } for s in subs]


@router.post("/submissions/save")
async def save_submission(
    period_id: int = Form(...),
    section_key: str = Form("__full__"),
    field_values: str = Form(...),
    status: str = Form("draft"),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    any_auth(user)
    if status not in ("draft", "submitted"):
        raise HTTPException(400, "Invalid status")
    try:
        fv = json.loads(field_values)
    except:
        raise HTTPException(400, "field_values must be valid JSON")
    result = await db.execute(
        select(ReportSubmission).where(
            ReportSubmission.period_id == period_id,
            ReportSubmission.user_id == user.id,
        )
    )
    existing = result.scalar_one_or_none()
    was_submitted = existing and existing.status != "submitted" and status == "submitted"
    if existing:
        existing.field_values = field_values
        existing.status = status
        if status == "submitted":
            existing.submitted_at = datetime.utcnow()
    else:
        sub = ReportSubmission(
            period_id=period_id,
            section_key="__full__",
            user_id=user.id,
            field_values=field_values,
            status=status,
            submitted_at=datetime.utcnow() if status == "submitted" else None,
        )
        db.add(sub)
        was_submitted = status == "submitted"
    if was_submitted:
        result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
        period = result.scalar_one_or_none()
        tname = "Unknown"
        plabel = "Unknown"
        if period:
            result = await db.execute(select(ReportTemplate).where(ReportTemplate.id == period.template_id))
            template = result.scalar_one_or_none()
            if template: tname = template.name
            plabel = period.label
        result = await db.execute(
            select(User).options(selectinload(User.role)).where(
                User.is_active == True,
                User.id != user.id,
            )
        )
        notify_users = [u for u in result.scalars().all() if u.role and u.role.name in ("admin", "ops_manager")]
        for nu in notify_users:
            db.add(Notification(
                user_id=nu.id,
                message=f'Report submitted: {user.name} filled "{tname}" — {plabel}',
                is_read=False,
            ))
    await db.commit()
    return {"msg": "Submission saved", "status": status}


# ─── Fill Status (Google Form style) ───

@router.get("/periods/{period_id}/fill-status")
async def get_fill_status(period_id: int, db: AsyncSession = Depends(get_db), user: User = Depends(get_current_user)):
    admin_or_ops(user)
    result = await db.execute(select(ReportPeriod).where(ReportPeriod.id == period_id))
    period = result.scalar_one_or_none()
    if not period:
        raise HTTPException(404, "Period not found")
    result = await db.execute(select(ReportTemplate).where(ReportTemplate.id == period.template_id))
    template = result.scalar_one_or_none()
    if not template:
        raise HTTPException(404, "Template not found")
    assigned_roles = json.loads(template.assigned_roles) if isinstance(template.assigned_roles, str) else (template.assigned_roles or [])
    result = await db.execute(select(User).options(selectinload(User.role)))
    all_users = result.scalars().all()
    target_users = [u for u in all_users if u.cpf in assigned_roles] if assigned_roles else all_users
    result = await db.execute(
        select(ReportSubmission).where(
            ReportSubmission.period_id == period_id,
            ReportSubmission.section_key == "__full__",
        )
    )
    submissions = result.scalars().all()
    submitted_user_ids = {s.user_id for s in submissions if s.user_id}
    fill_status = []
    for u in target_users:
        has_filled = u.id in submitted_user_ids
        sub = next((s for s in submissions if s.user_id == u.id), None)
        fill_status.append({
            "user_id": u.id,
            "cpf": u.cpf,
            "name": u.name,
            "role": u.role.name if u.role else None,
            "has_filled": has_filled,
            "status": sub.status if sub else None,
            "submitted_at": str(sub.submitted_at) if sub and sub.submitted_at else None,
        })
    return {"template_name": template.name, "period_label": period.label, "assigned_roles": assigned_roles, "fill_status": fill_status}

@router.get("/export/{period_id}")
async def export_report(
    period_id: int,
    format: str = "json",
    db: AsyncSession = Depends(get_db),
    user: User = Depends(get_current_user),
):
    any_auth(user)
    result = await db.execute(
        select(ReportPeriod).options(selectinload(ReportPeriod.template)).where(ReportPeriod.id == period_id)
    )
    period = result.scalar_one_or_none()
    if not period:
        raise HTTPException(404, "Period not found")

    template_sections = json.loads(period.template.sections) if isinstance(period.template.sections, str) else period.template.sections

    result = await db.execute(
        select(ReportSubmission).options(selectinload(ReportSubmission.assignee), selectinload(ReportSubmission.submitter)).where(ReportSubmission.period_id == period_id)
    )
    subs = result.scalars().all()

    submission_map = {}
    for s in subs:
        if s.section_key == "__full__":
            submission_map["__full__"] = {
                "field_values": json.loads(s.field_values) if isinstance(s.field_values, str) else s.field_values,
                "status": s.status,
                "user_name": s.submitter.name if s.submitter else "—",
                "user_cpf": s.submitter.cpf if s.submitter else "—",
                "submitted_at": str(s.submitted_at) if s.submitted_at else "—",
            }
        else:
            submission_map[s.section_key] = {
                "field_values": json.loads(s.field_values) if isinstance(s.field_values, str) else s.field_values,
                "status": s.status,
                "assigned_to_name": s.assignee.name if s.assignee else "Unassigned",
                "submitted_at": str(s.submitted_at) if s.submitted_at else "—",
            }

    period_assignments = json.loads(period.section_assignments) if period.section_assignments else {}

    # Build user-id → name mapping
    result = await db.execute(select(User))
    all_users = result.scalars().all()
    user_names = {u.id: u.name for u in all_users}

    report_data = {
        "template_name": period.template.name,
        "period_label": period.label,
        "period_start": str(period.start_date) if period.start_date else None,
        "period_end": str(period.end_date) if period.end_date else None,
        "sections": [],
    }
    for sec in template_sections:
        sec_key = sec.get("key", "")
        sub_data = submission_map.get(sec_key, {})
        assigned_uid = period_assignments.get(sec_key) or sub_data.get("assigned_to")
        assigned_name = user_names.get(assigned_uid) if assigned_uid else sub_data.get("assigned_to_name", "Unassigned")
        report_data["sections"].append({
            "title": sec.get("title", sec_key),
            "key": sec_key,
            "fields": sec.get("fields", []),
            "values": sub_data.get("field_values", {}),
            "status": sub_data.get("status", "pending"),
            "assigned_to": assigned_name,
            "submitted_at": sub_data.get("submitted_at", "—"),
        })

    if format == "json":
        return report_data

    elif format == "text":
        lines = []
        lines.append(f"{'='*60}")
        lines.append(f"  {report_data['template_name']}")
        lines.append(f"  Period: {report_data['period_label']}")
        lines.append(f"{'='*60}\n")
        for sec in report_data["sections"]:
            lines.append(f"── {sec['title']} ──")
            lines.append(f"   Assigned to: {sec['assigned_to']}")
            lines.append(f"   Status: {sec['status']}")
            for f in sec["fields"]:
                fk = f.get("key", "")
                fl = f.get("label", fk)
                val = sec["values"].get(fk, "—")
                lines.append(f"   {fl}: {val}")
            lines.append("")
        return Response("\n".join(lines), media_type="text/plain")

    elif format == "html":
        rows = ""
        for sec in report_data["sections"]:
            fields_html = ""
            for f in sec["fields"]:
                fk = f.get("key", "")
                fl = f.get("label", fk)
                val = sec["values"].get(fk, "—")
                fields_html += f"<tr><td style='padding:4px 8px;color:#555;font-weight:600'>{fl}</td><td style='padding:4px 8px'>{val}</td></tr>"
            status_color = {"draft":"#E65100","submitted":"#1565C0","pending":"#999"}
            rows += f"""
            <div style="margin-bottom:16px;border:1px solid #e0e4e8;border-radius:8px;padding:12px">
                <div style="font-size:14px;font-weight:700;color:#0b3d91;margin-bottom:4px">{sec['title']}</div>
                <div style="font-size:12px;color:#888;margin-bottom:8px">Assigned: {sec['assigned_to']} | Status: <span style="color:{status_color.get(sec['status'],'#999')}">{sec['status']}</span></div>
                <table style="width:100%;border-collapse:collapse;font-size:13px">{fields_html}</table>
            </div>"""
        html = f"""<html><head><meta charset="utf-8"><title>{report_data['template_name']}</title></head><body style="font-family:sans-serif;padding:20px;max-width:800px;margin:auto">
            <h1 style="color:#0b3d91">{report_data['template_name']}</h1>
            <p style="color:#666">Period: {report_data['period_label']}</p>
            {rows}
        </body></html>"""
        return Response(html, media_type="text/html")

    elif format == "docx":
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            doc = Document()
            doc.add_heading(report_data["template_name"], level=1)
            doc.add_paragraph(f"Period: {report_data['period_label']}")
            for sec in report_data["sections"]:
                doc.add_heading(sec["title"], level=2)
                doc.add_paragraph(f"Assigned to: {sec['assigned_to']}  |  Status: {sec['status']}")
                table = doc.add_table(rows=1, cols=2)
                table.style = "Light Grid Accent 1"
                hdr = table.rows[0].cells
                hdr[0].text = "Field"
                hdr[1].text = "Value"
                for f in sec["fields"]:
                    fk = f.get("key", "")
                    fl = f.get("label", fk)
                    val = sec["values"].get(fk, "—")
                    row = table.add_row().cells
                    row[0].text = fl
                    row[1].text = str(val)
                doc.add_paragraph("")
            import io
            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)
            return Response(
                buf.read(),
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f"attachment; filename={period.template.name.replace(' ','_')}_{period.label.replace(' ','_')}.docx"}
            )
        except ImportError:
            raise HTTPException(500, "DOCX export not available (python-docx not installed)")

    elif format == "pptx":
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = report_data["template_name"]
            slide.placeholders[1].text = f"Period: {report_data['period_label']}"
            for sec in report_data["sections"]:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = sec["title"]
                txBox = slide.placeholders[1]
                tf = txBox.text_frame
                tf.text = f"Assigned to: {sec['assigned_to']}\nStatus: {sec['status']}\n"
                for f in sec["fields"]:
                    fk = f.get("key", "")
                    fl = f.get("label", fk)
                    val = sec["values"].get(fk, "—")
                    p = tf.add_paragraph()
                    p.text = f"{fl}: {val}"
                    p.space_after = Pt(4)
            import io
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return Response(
                buf.read(),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename={period.template.name.replace(' ','_')}_{period.label.replace(' ','_')}.pptx"}
            )
        except ImportError:
            raise HTTPException(500, "PPTX export not available (python-pptx not installed)")

    elif format == "pdf":
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            import io
            buf = io.BytesIO()
            doc = SimpleDocTemplate(buf, pagesize=A4)
            styles = getSampleStyleSheet()
            elements = []
            elements.append(Paragraph(report_data["template_name"], styles["Title"]))
            elements.append(Paragraph(f"Period: {report_data['period_label']}", styles["Normal"]))
            elements.append(Spacer(1, 12))
            for sec in report_data["sections"]:
                elements.append(Paragraph(sec["title"], styles["Heading2"]))
                elements.append(Paragraph(f"Assigned to: {sec['assigned_to']} | Status: {sec['status']}", styles["Normal"]))
                elements.append(Spacer(1, 6))
                data = [["Field", "Value"]]
                for f in sec["fields"]:
                    fk = f.get("key", "")
                    fl = f.get("label", fk)
                    val = sec["values"].get(fk, "—")
                    data.append([fl, str(val)])
                tbl = Table(data, colWidths=[200, 350])
                tbl.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0b3d91")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f7fa")]),
                ]))
                elements.append(tbl)
                elements.append(Spacer(1, 12))
            doc.build(elements)
            buf.seek(0)
            return Response(
                buf.read(),
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename={period.template.name.replace(' ','_')}_{period.label.replace(' ','_')}.pdf"}
            )
        except ImportError:
            raise HTTPException(500, "PDF export not available (reportlab not installed)")

    raise HTTPException(400, f"Unsupported format: {format}")
